import streamlit as st
import fitz  # PyMuPDF for handling PDFs
from pptx import Presentation
import tempfile
import os
import requests  # Import requests for making HTTP POST requests

# # Function to process the uploaded file and extract text
# def process_file(file_path):
#     _, file_extension = os.path.splitext(file_path)
#     if file_extension == '.pdf':
#         return extract_text_from_pdf(file_path)
#     elif file_extension == '.pptx':  # Ensure correct extension is used
#         return extract_text_from_ppt(file_path)
#     else:
#         return "Unsupported file type"

# # Function to extract text from PDF
# def extract_text_from_pdf(file_path):
#     doc = fitz.open(file_path)
#     text = ""
#     for page in doc:
#         text += page.get_text()
#     return text

# # Function to extract text from PPTX
# def extract_text_from_ppt(file_path):
#     prs = Presentation(file_path)
#     text = ""
#     for slide in prs.slides:
#         for shape in slide.shapes:
#             if hasattr(shape, "text"):
#                 text += shape.text
#     return text

# # Function to make POST request to the FastAPI server to save documents to Pinecone
# def save_to_pinecone(extracted_texts):
#     try:
#         response = requests.post("http://localhost:8000/embeed", json=extracted_texts)
#         if response.status_code == 200:
#             st.success("Documents successfully saved to Pinecone!")
#         else:
#             st.error(f"Failed to save documents to Pinecone. Status code: {response.status_code}")
#     except Exception as e:
#         st.error(f"An error occurred: {e}")

# Function to make POST request to the FastAPI server to analyze the question
def analyze_question(question):
    try:
        response = requests.post("http://localhost:8000/analyze", json={"question": question})
        if response.status_code == 200:
            answer = response.json().get("answer", "No answer found.")
            st.success(f"Answer: {answer}")
        else:
            st.error(f"Failed to analyze the question. Status code: {response.status_code}")
    except Exception as e:
        st.error(f"An error occurred: {e}")

# Function to display the login page
def login_page():
    st.title("Document Analyzer - Login")
    st.subheader("Please log in")

    username = st.text_input("Username")
    password = st.text_input("Password", type="password")

    if st.button("Login"):
        if username == 'Ricky' and password == 'Aipersona2024':
            st.session_state.logged_in = True
            st.success("Login successful!")
        else:
            st.error("Incorrect username or password")

# Main Streamlit app logic
def main_page():
    st.title("Document Analyzer")

    # st.subheader("Upload PDF or PPTX files")
    # uploaded_files = st.file_uploader("Choose files", type=['pdf', 'pptx'], accept_multiple_files=True)

    # if uploaded_files:
    #     temp_dir = tempfile.mkdtemp()
    #     extracted_texts = {}

    #     for uploaded_file in uploaded_files:
    #         file_path = os.path.join(temp_dir, uploaded_file.name)

    #         with open(file_path, "wb") as f:
    #             f.write(uploaded_file.getbuffer())

    #         # Extract text from the file
    #         extracted_text = process_file(file_path)
    #         extracted_texts[uploaded_file.name] = extracted_text

    #         # Display extracted text
    #         st.text_area(f"Extracted Text from {uploaded_file.name}", extracted_text, height=200)

        # # Button to save extracted texts to Pinecone via the FastAPI endpoint
        # if st.button("Save to Pinecone"):
        #     with st.spinner("Saving to Pinecone..."):
        #         save_to_pinecone(extracted_texts)

    # Analyze Text Section
    st.subheader("Analyze Text")
    st.code("What to do if a customer tells you they plan to purchase new Cisco gear to replace older Cisco or competitor hardware? Write in 5 steps how to purchase new Cisco Gear")
    question = st.text_input("Enter your question")

    if st.button("Analyze"):
        with st.spinner("Analyzing..."):
            analyze_question(question)

# Main application logic
def run():
    if 'logged_in' not in st.session_state:
        login_page()
    else:
        main_page()
